from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
import os
import PyPDF2
import docx
from pptx import Presentation
import io
import re
from typing import List
import logging
import openai

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configure for Azure OpenAI
openai.api_type = "azure"
openai.api_key = os.getenv("AZURE_OPENAI_KEY")
openai.api_base = os.getenv("AZURE_OPENAI_ENDPOINT")
openai.api_version = "2024-02-01"

app = FastAPI(
    title="Tacit Knowledge Bot",
    description="AI-powered document knowledge assistant",
    version="1.0.0"
)

# Global storage for documents
documents = []

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Serve static files
try:
    app.mount("/static", StaticFiles(directory="static"), name="static")
except:
    pass

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        return ""

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from Word document"""
    try:
        doc_file = io.BytesIO(file_content)
        doc = docx.Document(doc_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting DOCX text: {e}")
        return ""

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint presentation"""
    try:
        ppt_file = io.BytesIO(file_content)
        prs = Presentation(ppt_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PPTX text: {e}")
        return ""

def find_relevant_documents(question: str, documents: List[dict], threshold: int = 1) -> List[dict]:
    """Find documents relevant to the question"""
    if not documents:
        return []
    
    stop_words = {'what', 'how', 'where', 'when', 'why', 'who', 'is', 'are', 'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'about', 'can', 'could', 'should', 'would', 'do', 'does', 'did'}
    
    question_words = [word.lower().strip('.,!?') for word in question.split() 
                     if len(word) > 2 and word.lower() not in stop_words]
    
    relevant_docs = []
    for doc in documents:
        content_lower = doc["content"].lower()
        relevance_score = sum(1 for word in question_words if word in content_lower)
        if relevance_score >= threshold:
            relevant_docs.append((doc, relevance_score))
    
    relevant_docs.sort(key=lambda x: x[1], reverse=True)
    return [doc[0] for doc in relevant_docs[:3]]

@app.get("/")
def read_root():
    """Serve the web interface"""
    return FileResponse('static/index.html')

@app.get("/api")
def api_status():
    """API status endpoint"""
    return {
        "message": "🤖 Tacit Knowledge Bot is LIVE!",
        "status": "healthy",
        "documents_loaded": len(documents),
        "azure_openai_available": bool(openai.api_key and openai.api_base)
    }

@app.get("/health")
def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "documents": len(documents),
        "azure_openai_connected": bool(openai.api_key and openai.api_base),
        "version": "1.0.0"
    }

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """Upload and process documents"""
    try:
        logger.info(f"Processing file: {file.filename}")
        
        filename = file.filename.lower()
        content = await file.read()
        
        # Extract text based on file type
        if filename.endswith('.txt'):
            text_content = content.decode("utf-8")
        elif filename.endswith('.pdf'):
            text_content = extract_text_from_pdf(content)
        elif filename.endswith('.docx') or filename.endswith('.doc'):
            text_content = extract_text_from_docx(content)
        elif filename.endswith('.pptx') or filename.endswith('.ppt'):
            text_content = extract_text_from_pptx(content)
        else:
            return {
                "status": "error", 
                "message": f"Unsupported file type: {file.filename}. Supports: PDF, Word, PowerPoint, TXT"
            }
        
        if not text_content.strip():
            return {
                "status": "error",
                "message": f"No text could be extracted from {file.filename}"
            }
        
        # Store document
        documents.append({
            "content": text_content,
            "filename": file.filename,
            "file_type": filename.split('.')[-1],
            "size": len(text_content)
        })
        
        logger.info(f"Successfully processed {file.filename}, extracted {len(text_content)} characters")
        
        return {
            "filename": file.filename,
            "status": "success",
            "message": f"Successfully processed {file.filename}! Total docs: {len(documents)}",
            "extracted_characters": len(text_content)
        }
        
    except Exception as e:
        logger.error(f"Error processing file {file.filename}: {e}")
        return {"status": "error", "message": f"Error processing file: {str(e)}"}

@app.post("/chat")
async def chat(question: str = Form(...)):
    """Chat with the knowledge bot"""
    try:
        if not openai.api_key or not openai.api_base:
            return {
                "question": question,
                "answer": "Azure OpenAI API key or endpoint not configured. Please check environment variables.",
                "source": "Error"
            }
        
        logger.info(f"Processing question: {question}")
        
        # Find relevant documents
        relevant_docs = find_relevant_documents(question, documents, threshold=1)
        
        if relevant_docs:
            context_parts = []
            doc_names = []
            for doc in relevant_docs:
                doc_type = "📧" if doc['file_type'] == 'email' else "📄"
                context_parts.append(f"=== {doc_type} {doc['filename']} ===\n{doc['content']}")
                doc_names.append(doc['filename'])
            
            context = "\n\n".join(context_parts)
            prompt = f"Based on these documents, answer the question clearly and concisely:\n\n{context}\n\nQuestion: {question}\n\nAnswer:"
            source_info = f" (Based on: {', '.join(doc_names)})"
        else:
            prompt = f"Answer this general question: {question}"
            source_info = " (General knowledge)"
        
        # Get AI response using Azure OpenAI API
        response = openai.ChatCompletion.create(
            engine="gpt-35-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful knowledge assistant. Be clear, concise, and cite sources when using document information."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=500,
            temperature=0.7
        )
        
        answer = response.choices[0].message.content
        
        logger.info(f"Successfully answered question using {len(relevant_docs)} documents")
        
        return {
            "question": question,
            "answer": answer + source_info,
            "documents_used": len(relevant_docs),
            "source_documents": [doc['filename'] for doc in relevant_docs] if relevant_docs else []
        }
        
    except Exception as e:
        logger.error(f"Error processing chat: {e}")
        return {"question": question, "answer": f"Error: {str(e)}"}

@app.get("/documents")
def list_documents():
    """List all uploaded documents"""
    return {
        "total_documents": len(documents),
        "documents": [
            {
                "filename": doc["filename"],
                "type": doc["file_type"],
                "size": doc.get("size", 0),
                "icon": "📧" if doc["file_type"] == "email" else "📄"
            } for doc in documents
        ]
    }

@app.delete("/documents")
def clear_documents():
    """Clear all documents"""
    global documents
    count = len(documents)
    documents = []
    return {"message": f"Cleared {count} documents", "remaining": 0}

if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 80))
    uvicorn.run(app, host="0.0.0.0", port=port)
