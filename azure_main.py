from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from openai import AzureOpenAI
import os
from dotenv import load_dotenv
import PyPDF2
import docx
from pptx import Presentation
import io
import re

load_dotenv()

app = FastAPI(title="Tacit Knowledge Bot")

# Initialize Azure OpenAI
client = AzureOpenAI(
    api_key=os.getenv("AZURE_OPENAI_KEY"),
    api_version="2024-02-01",
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT")
)

documents = []

app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

# Mount static files
try:
    app.mount("/static", StaticFiles(directory="static"), name="static")
except:
    pass  # Skip if static folder doesn't exist

def extract_text_from_pdf(file_content):
    pdf_file = io.BytesIO(file_content)
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file_content):
    doc_file = io.BytesIO(file_content)
    doc = docx.Document(doc_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(file_content):
    ppt_file = io.BytesIO(file_content)
    prs = Presentation(ppt_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def find_relevant_documents(question, documents, threshold=1):
    if not documents:
        return []
    
    stop_words = {'what', 'how', 'where', 'when', 'why', 'who', 'is', 'are', 'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'about'}
    question_words = [word.lower().strip('.,!?') for word in question.split() 
                     if len(word) > 2 and word.lower() not in stop_words]
    
    relevant_docs = []
    for doc in documents:
        content_lower = doc["content"].lower()
        relevance_score = sum(1 for word in question_words if word in content_lower)
        if relevance_score >= threshold:
            relevant_docs.append((doc, relevance_score))
    
    relevant_docs.sort(key=lambda x: x[1], reverse=True)
    return [doc[0] for doc in relevant_docs[:3]]

@app.get("/")
def read_root():
    return {"message": "🤖 Tacit Knowledge Bot is running on Azure!", "status": "healthy"}

@app.get("/health")
def health_check():
    return {"status": "healthy", "message": "Bot is running!", "documents": len(documents)}

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    try:
        filename = file.filename.lower()
        content = await file.read()
        
        if filename.endswith('.txt'):
            text_content = content.decode("utf-8")
        elif filename.endswith('.pdf'):
            text_content = extract_text_from_pdf(content)
        elif filename.endswith('.docx') or filename.endswith('.doc'):
            text_content = extract_text_from_docx(content)
        elif filename.endswith('.pptx') or filename.endswith('.ppt'):
            text_content = extract_text_from_pptx(content)
        else:
            return {"status": "error", "message": f"Unsupported file type: {file.filename}"}
        
        if not text_content.strip():
            return {"status": "error", "message": f"No text extracted from {file.filename}"}
        
        documents.append({
            "content": text_content,
            "filename": file.filename,
            "file_type": filename.split('.')[-1]
        })
        
        return {
            "filename": file.filename,
            "status": "success",
            "message": f"✅ Successfully processed {file.filename}! Total docs: {len(documents)}"
        }
        
    except Exception as e:
        return {"status": "error", "message": f"Error: {str(e)}"}

@app.post("/chat")
async def chat(question: str = Form(...)):
    try:
        relevant_docs = find_relevant_documents(question, documents, threshold=1)
        
        if relevant_docs:
            context_parts = []
            doc_names = []
            for doc in relevant_docs:
                context_parts.append(f"=== {doc['filename']} ===\n{doc['content']}")
                doc_names.append(doc['filename'])
            
            context = "\n\n".join(context_parts)
            prompt = f"Based on these documents, answer the question:\n\n{context}\n\nQuestion: {question}\n\nAnswer:"
            source_info = f" (Based on: {', '.join(doc_names)})"
        else:
            prompt = f"Answer this question: {question}"
            source_info = " (General knowledge)"
        
        response = client.chat.completions.create(
            model="gpt-35-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant. Be clear and concise."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=500
        )
        
        return {
            "question": question, 
            "answer": response.choices[0].message.content + source_info,
            "documents_used": len(relevant_docs)
        }
        
    except Exception as e:
        return {"question": question, "answer": f"Error: {str(e)}"}

@app.get("/documents")
def list_documents():
    return {
        "total_documents": len(documents),
        "documents": [{"filename": doc["filename"], "type": doc["file_type"]} for doc in documents]
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=int(os.getenv("PORT", 80)))
